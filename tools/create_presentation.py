from pptx import Presentation
from pptx.util import Inches, Pt

slides = [
    {
        'title': 'Cyber Watch — Emotion-aware Threat Detection',
        'bullets': ['Multimodal detection (Text | Face | Voice)', 'Prototype demo: webcam + microphone + files'],
        'notes': 'Introduce the project and the demo goal: a prototype that detects threats in text, voice and facial expressions to augment security monitoring.'
    },
    {
        'title': 'Problem Statement',
        'bullets': [
            'Multimodal threats: phishing messages, aggressive voice calls, suspicious video behavior',
            'Current tools miss emotional/contextual cues',
            'Human moderators face high-volume multimedia streams'
        ],
        'notes': 'Give examples and emphasize the need to fuse signals (text, voice, face) to improve detection and prioritization.'
    },
    {
        'title': 'Objectives',
        'bullets': [
            'Modular prototype analyzing text, voice, and facial emotion',
            'Usable GUI for analysts with logging',
            'Lightweight fallbacks + support for heavy models',
            'Real-time monitoring and alerting'
        ],
        'notes': 'Clarify modularity and practical constraints (optional heavy model downloads).'
    },
    {
        'title': 'Methodology',
        'bullets': [
            'Text: RoBERTa + heuristics',
            'Voice: MFCC/spectral features + transformer ASR/emotion (fallback: RandomForest)',
            'Face: YOLO / DeepFace / Haar cascade fallbacks',
            'UI: Tkinter dashboard with analyzers and live monitoring'
        ],
        'notes': 'Explain cascaded/fallback approach and voice pipeline.'
    },
    {
        'title': 'System Design (Architecture)',
        'bullets': [
            'Frontend: main UI and modular GUIs',
            'Models: voice_model.py, text_model.py',
            'Utilities: file utils, DB history',
            'Optional models directory for heavy weights'
        ],
        'notes': 'Walk through data flow: input -> features -> models -> alert/log.'
    },
    {
        'title': 'Implementation Highlights',
        'bullets': [
            'Robust GUI with theme, tooltips, back navigation',
            'VoiceThreatClassifier: transformer pipelines or RandomForest fallback',
            'FacialEmotionAnalyzer: YOLO/DeepFace/Haar cascades',
            'Text classifier + Gmail scanner integration (optional)'
        ],
        'notes': 'Point to key files (main.py, gui/voice_gui.py, model/*) and defensive imports.'
    },
    {
        'title': 'Demo & Results',
        'bullets': [
            'Planned demo steps: voice file analysis, webcam face analyzer, text threat popup',
            'Results: file-based voice/text tests passed; live webcam validated locally',
            'Alerts: popup + beep + log entries'
        ],
        'notes': 'Outline demo flow and mention limitations (model downloads, audio devices).'
    },
    {
        'title': 'Conclusion & Next Steps',
        'bullets': [
            'Prototype shows multimodal detection and real-time alerts',
            'Next: CI, packaged demo weights, labeled dataset & evaluation',
            'UX polish and dockerized deployment'
        ],
        'notes': 'Summarize wins and next work items; ask for sample data or resources for large-model demos.'
    }
]

prs = Presentation()
# Use a simple title slide layout for the first slide
for i, s in enumerate(slides):
    if i == 0:
        slide_layout = prs.slide_layouts[0]  # title slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = s['title']
        subtitle.text = s['bullets'][0] if s['bullets'] else ''
        # add notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = s['notes']
    else:
        slide_layout = prs.slide_layouts[1]  # title and content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = s['title']
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for b in s['bullets']:
            p = body.add_paragraph()
            p.text = b
            p.level = 0
            p.font = None
        # add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = s['notes']

out_path = 'CyberWatch_presentation.pptx'
prs.save(out_path)
print(f'Created presentation: {out_path}')
